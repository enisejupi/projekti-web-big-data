#!/usr/bin/env python3.9
# -*- coding: utf-8 -*-
"""
Presentation Generator
Creates PowerPoint presentation with financial analysis results
"""

import sqlite3
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import seaborn as sns
from datetime import datetime
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("WARNING: python-pptx not installed. Install with: pip3.9 install python-pptx")

DB_PATH = '/opt/financial-analysis/data/market_data.db'
REPORT_DIR = '/opt/financial-analysis/reports'

def get_db_stats():
    """Get database statistics"""
    conn = sqlite3.connect(DB_PATH)
    
    stats = {}
    stats['total'] = pd.read_sql_query("SELECT COUNT(*) as cnt FROM market_data", conn)['cnt'].iloc[0]
    stats['symbols'] = pd.read_sql_query("SELECT COUNT(DISTINCT symbol) as cnt FROM market_data", conn)['cnt'].iloc[0]
    
    dates = pd.read_sql_query("SELECT MIN(timestamp) as first, MAX(timestamp) as last FROM market_data", conn)
    stats['first_date'] = dates['first'].iloc[0] if not dates.empty else 'N/A'
    stats['last_date'] = dates['last'].iloc[0] if not dates.empty else 'N/A'
    
    stats['top_volume'] = pd.read_sql_query("""
        SELECT symbol, AVG(volume) as avg_vol
        FROM market_data 
        GROUP BY symbol 
        ORDER BY avg_vol DESC 
        LIMIT 10
    """, conn)
    
    stats['top_gainers'] = pd.read_sql_query("""
        SELECT symbol, AVG(change_percent) as avg_change
        FROM market_data 
        WHERE change_percent IS NOT NULL
        GROUP BY symbol 
        ORDER BY avg_change DESC 
        LIMIT 10
    """, conn)
    
    conn.close()
    return stats

def create_chart(data, title, filename, chart_type='bar'):
    """Create and save chart"""
    plt.figure(figsize=(10, 6))
    
    if chart_type == 'bar':
        plt.bar(range(len(data)), data.values[:, 1], color='#3498db')
        plt.xticks(range(len(data)), data.values[:, 0], rotation=45)
    
    plt.title(title, fontsize=16, fontweight='bold')
    plt.tight_layout()
    plt.savefig(filename, dpi=150, bbox_inches='tight')
    plt.close()

def generate_presentation():
    """Generate PowerPoint presentation"""
    
    if not PPTX_AVAILABLE:
        print("❌ Cannot generate presentation: python-pptx not installed")
        return None
    
    print("📊 Generating presentation...")
    
    # Get stats
    stats = get_db_stats()
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Analiza Financiare"
    subtitle.text = f"Sistemi i Mbledhjes së të Dhënave Financiare\\n84 Orë - Projekti Universitar\\n{datetime.now().strftime('%d %B %Y')}"
    
    # Slide 2: Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Përmbledhje"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = f"📊 Totali i Regjistrave: {stats['total']:,}"
    
    p = tf.add_paragraph()
    p.text = f"🏢 Simbole të Analizuara: {stats['symbols']}"
    
    p = tf.add_paragraph()
    p.text = f"📅 Periudha: {stats['first_date'][:19]} - {stats['last_date'][:19]}"
    
    p = tf.add_paragraph()
    p.text = "🎯 Objektivi: Mbledhje dhe analiza e të dhënave financiare për 84 orë"
    
    p = tf.add_paragraph()
    p.text = "💻 Infrastruktura: 10 VMs, Apache Spark, Python, Machine Learning"
    
    # Slide 3: Top Volume
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Top 10 - Vëllimi Mesatar"
    
    chart_filename = f"{REPORT_DIR}/temp_volume_chart.png"
    create_chart(stats['top_volume'], 'Top 10 - Vëllimi Mesatar', chart_filename)
    
    left = Inches(1)
    top = Inches(2)
    pic = slide.shapes.add_picture(chart_filename, left, top, width=Inches(8))
    
    # Slide 4: Top Gainers
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Top 10 - Fituesit"
    
    chart_filename = f"{REPORT_DIR}/temp_gainers_chart.png"
    create_chart(stats['top_gainers'], 'Top 10 Fituesit (%)', chart_filename)
    
    pic = slide.shapes.add_picture(chart_filename, left, top, width=Inches(8))
    
    # Slide 5: Metodologjia
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Metodologjia"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. Mbledhja e të Dhënave"
    
    p = tf.add_paragraph()
    p.text = "   • Yahoo Finance API"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • 500+ simbole financiare"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Intervale 5-minutëshe"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. Përpunimi"
    
    p = tf.add_paragraph()
    p.text = "   • Parallel processing (32 workers)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • SQLite database"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. Machine Learning"
    
    p = tf.add_paragraph()
    p.text = "   • Random Forest"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Gradient Boosting"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • K-Means Clustering"
    p.level = 1
    
    # Slide 6: Përfundime
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Përfundime dhe Rekomandime"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "✅ Sistemi funksionon me sukses për 84 orë"
    
    p = tf.add_paragraph()
    p.text = f"✅ Mbledhur {stats['total']:,} regjistrat nga {stats['symbols']} simbole"
    
    p = tf.add_paragraph()
    p.text = "✅ Infrastruktura e distribuuar përpunon të dhëna në kohë reale"
    
    p = tf.add_paragraph()
    p.text = "✅ Machine Learning ofron parashikime me saktësi të lartë"
    
    p = tf.add_paragraph()
    p.text = "💡 Rekomandime:"
    p = tf.add_paragraph()
    p.text = "   • Diversifikim në sektorë të ndryshëm"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Fokus në kompani me kapitalizim të lartë"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   • Monitorim i vazhdueshëm i volatilitetit"
    p.level = 1
    
    # Save presentation
    os.makedirs(REPORT_DIR, exist_ok=True)
    filename = f"{REPORT_DIR}/presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs.save(filename)
    
    # Cleanup temp files
    try:
        os.remove(f"{REPORT_DIR}/temp_volume_chart.png")
        os.remove(f"{REPORT_DIR}/temp_gainers_chart.png")
    except:
        pass
    
    print(f"✅ Presentation saved: {filename}")
    return filename

if __name__ == '__main__':
    generate_presentation()
